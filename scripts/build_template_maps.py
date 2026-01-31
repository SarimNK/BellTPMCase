import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

BASE = Path(__file__).resolve().parents[1]
RESOURCES = BASE / 'src' / 'resources'

TEMPLATES = {
    'bell-enterprise': RESOURCES / 'BellCanadaTemplate.pptx',
    'internal-strategy': RESOURCES / '11b - alien minds - slides.pptx',
    'quarterly-review': RESOURCES / '1%2B-%2BIntroduction%2Bto%2B2244.pptx',
}

PLACEHOLDER_TYPE_MAP = {
    PP_PLACEHOLDER.TITLE: 'title',
    PP_PLACEHOLDER.CENTER_TITLE: 'title',
    PP_PLACEHOLDER.SUBTITLE: 'subtitle',
    PP_PLACEHOLDER.BODY: 'body',
    PP_PLACEHOLDER.OBJECT: 'body',
}

output = {}

for template_id, path in TEMPLATES.items():
    if not path.exists():
        raise FileNotFoundError(f"Missing template: {path}")
    prs = Presentation(str(path))
    layouts = []

    for i, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            ptype = PLACEHOLDER_TYPE_MAP.get(shape.placeholder_format.type)
            placeholders.append({
                'name': shape.name,
                'idx': shape.placeholder_format.idx,
                'type': ptype or 'other'
            })
        layouts.append({
            'index': i,
            'name': layout.name,
            'placeholders': placeholders
        })

    def find_layout_index(require_title=True, require_body=False, require_subtitle=False):
        for layout in layouts:
            types = {p['type'] for p in layout['placeholders']}
            if require_title and 'title' not in types:
                continue
            if require_body and 'body' not in types:
                continue
            if require_subtitle and 'subtitle' not in types:
                continue
            return layout['index']
        return 0

    def get_placeholder_idx(layout_index, ptype):
        layout = next((l for l in layouts if l['index'] == layout_index), None)
        if not layout:
            return None
        for p in layout['placeholders']:
            if p['type'] == ptype:
                return p['idx']
        return None

    title_layout_index = find_layout_index(require_title=True, require_body=False)
    content_layout_index = find_layout_index(require_title=True, require_body=True)

    output[template_id] = {
        'template_file': str(path.name),
        'title_layout_index': title_layout_index,
        'content_layout_index': content_layout_index,
        'title_placeholder_idx': get_placeholder_idx(content_layout_index, 'title'),
        'body_placeholder_idx': get_placeholder_idx(content_layout_index, 'body'),
        'subtitle_placeholder_idx': get_placeholder_idx(title_layout_index, 'subtitle'),
        'layouts': layouts
    }

(RESOURCES / 'template-maps.json').write_text(json.dumps(output, indent=2))
print('Wrote template-maps.json')
