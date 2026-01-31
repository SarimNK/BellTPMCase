import io
import json
import os
import base64
from typing import List, Optional

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


def delete_all_slides(prs: Presentation):
    """Delete all existing slides from the presentation, keeping only layouts."""
    # We need to delete slides from the XML structure
    # The slides collection is in prs.slides._sldIdLst
    xml_slides = prs.slides._sldIdLst
    slides_to_delete = list(xml_slides)
    for sldId in slides_to_delete:
        # Get the relationship ID and remove it
        rId = sldId.rId
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


class SlidePayload(BaseModel):
    title: Optional[str] = None
    bullets: List[str] = []
    layout_index: Optional[int] = None
    title_placeholder_idx: Optional[int] = None
    body_placeholder_idx: Optional[int] = None
    subtitle_placeholder_idx: Optional[int] = None
    subtitle: Optional[str] = None
    image_ref: Optional[str] = None
    image_data: Optional[str] = None  # Base64 encoded image
    image_position: Optional[str] = "right"  # 'left' or 'right'


class DeckPayload(BaseModel):
    title: Optional[str] = None
    slides: List[SlidePayload]


class ExportRequest(BaseModel):
    templateId: str
    internalUseOnly: bool = True
    deck: DeckPayload


def load_template_maps():
    base_dir = os.path.dirname(__file__)
    maps_path = os.path.abspath(os.path.join(base_dir, "..", "src", "resources", "template-maps.json"))
    if not os.path.exists(maps_path):
        raise FileNotFoundError(f"template-maps.json not found at {maps_path}")
    with open(maps_path, "r", encoding="utf-8") as f:
        return json.load(f)


def get_template_path(template_file: str):
    base_dir = os.path.dirname(__file__)
    return os.path.abspath(os.path.join(base_dir, "..", "src", "resources", template_file))


def safe_set_placeholder(slide, idx: Optional[int], text: Optional[str]):
    """Set text on a placeholder if it exists."""
    if idx is None or text is None:
        return False
    try:
        placeholder = slide.placeholders[idx]
        placeholder.text = text
        return True
    except (KeyError, IndexError):
        return False


def safe_set_body_with_bullets(slide, idx: Optional[int], bullets: List[str]):
    """Set body placeholder with proper bullet formatting."""
    if idx is None or not bullets:
        return False
    try:
        placeholder = slide.placeholders[idx]
        tf = placeholder.text_frame
        tf.clear()  # Clear existing content
        
        for i, bullet_text in enumerate(bullets):
            if not bullet_text:
                continue
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet_text
            p.level = 0  # Top-level bullet
        return True
    except (KeyError, IndexError):
        return False


def add_image_to_slide(slide, image_data: str, has_bullets: bool = True, position: str = "right"):
    """Add an image to the slide from base64 data."""
    if not image_data:
        return False
    
    try:
        # Remove data URL prefix if present
        if ',' in image_data:
            image_data = image_data.split(',', 1)[1]
        
        # Decode base64 to bytes
        image_bytes = base64.b64decode(image_data)
        image_stream = io.BytesIO(image_bytes)
        
        # Position image based on settings
        if has_bullets:
            if position == "left":
                # Image on the left side
                left = Inches(0.5)
            else:
                # Image on the right side (default)
                left = Inches(5.5)
            top = Inches(1.5)
            width = Inches(4.0)
            height = Inches(3.5)
        else:
            # Image centered and larger when no bullets
            left = Inches(2.0)
            top = Inches(1.5)
            width = Inches(6.0)
            height = Inches(4.5)
        
        # Add picture - it will maintain aspect ratio within the bounds
        slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
        return True
    except Exception as e:
        print(f"Error adding image: {e}")
        return False


def remove_watermarks(slide):
    """Remove watermark shapes from a slide (e.g., Harris Poll, copyright notices)."""
    # Keywords that indicate a watermark
    watermark_keywords = [
        'harris', 'poll', 'copyright', '©', 'confidential', 
        'proprietary', 'draft', 'sample', 'template',
        'all rights reserved', 'www.', 'http'
    ]
    
    # Shapes to remove (collect first, then remove to avoid iteration issues)
    shapes_to_remove = []
    
    for shape in slide.shapes:
        try:
            # Check if shape has text
            if hasattr(shape, 'text_frame'):
                text = shape.text_frame.text.lower().strip()
                if text:
                    # Check if it's in the bottom portion of the slide (footer area)
                    # Slide height is typically 7.5 inches, check if shape is in bottom 1.5 inches
                    shape_bottom = shape.top + shape.height
                    is_in_footer = shape.top > Inches(6.0)  # Bottom 1.5 inches
                    
                    # Check for watermark keywords
                    is_watermark = any(keyword in text for keyword in watermark_keywords)
                    
                    # Also remove if it's a small text box in footer with any text
                    is_small_footer_text = is_in_footer and shape.height < Inches(0.5)
                    
                    if is_watermark or (is_small_footer_text and len(text) < 100):
                        shapes_to_remove.append(shape)
            
            # Check for picture shapes that might be watermark logos
            elif hasattr(shape, 'shape_type'):
                # If it's a small image in the footer area, might be a logo watermark
                if shape.top > Inches(6.5) and shape.height < Inches(0.75):
                    shapes_to_remove.append(shape)
                    
        except Exception:
            continue
    
    # Remove the watermark shapes
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except Exception:
            pass


@app.get("/api/health")
def health():
    return {"status": "ok"}


@app.post("/api/pptx")
def generate_pptx(payload: ExportRequest):
    template_maps = load_template_maps()
    template_map = template_maps.get(payload.templateId)
    if not template_map:
        raise HTTPException(status_code=400, detail=f"Unknown templateId: {payload.templateId}")

    template_path = get_template_path(template_map["template_file"])
    if not os.path.exists(template_path):
        raise HTTPException(status_code=404, detail=f"Template file not found: {template_path}")

    prs = Presentation(template_path)
    
    # CRITICAL: Delete all existing slides from the template
    # We only want the layouts/masters, not the existing content
    delete_all_slides(prs)

    title_layout_index = template_map.get("title_layout_index", 0)
    content_layout_index = template_map.get("content_layout_index", 1)
    title_placeholder_idx = template_map.get("title_placeholder_idx", 0)
    body_placeholder_idx = template_map.get("body_placeholder_idx", 1)
    subtitle_placeholder_idx = template_map.get("subtitle_placeholder_idx")

    for i, slide_payload in enumerate(payload.deck.slides):
        is_title = i == 0
        layout_index = slide_payload.layout_index
        if layout_index is None:
            layout_index = title_layout_index if is_title else content_layout_index

        try:
            slide_layout = prs.slide_layouts[layout_index]
        except (IndexError, KeyError):
            # Fallback to content layout if specified layout doesn't exist
            try:
                slide_layout = prs.slide_layouts[content_layout_index]
            except (IndexError, KeyError):
                slide_layout = prs.slide_layouts[0]

        slide = prs.slides.add_slide(slide_layout)
        
        # Remove any watermarks inherited from the layout
        remove_watermarks(slide)

        # Set the title placeholder
        title_idx = slide_payload.title_placeholder_idx if slide_payload.title_placeholder_idx is not None else title_placeholder_idx
        safe_set_placeholder(slide, title_idx, slide_payload.title)

        # Set subtitle for title slides
        if is_title and slide_payload.subtitle:
            subtitle_idx = slide_payload.subtitle_placeholder_idx if slide_payload.subtitle_placeholder_idx is not None else subtitle_placeholder_idx
            safe_set_placeholder(slide, subtitle_idx, slide_payload.subtitle)

        # Set body/bullets with proper formatting
        if slide_payload.bullets:
            body_idx = slide_payload.body_placeholder_idx if slide_payload.body_placeholder_idx is not None else body_placeholder_idx
            # Try to set with bullet formatting first
            if not safe_set_body_with_bullets(slide, body_idx, slide_payload.bullets):
                # Fallback to simple text
                body_text = "\n".join([b for b in slide_payload.bullets if b])
                safe_set_placeholder(slide, body_idx, body_text)

        # Add image if present
        if slide_payload.image_data:
            has_bullets = bool(slide_payload.bullets)
            position = slide_payload.image_position or "right"
            add_image_to_slide(slide, slide_payload.image_data, has_bullets, position)

        # Add Bell Confidential watermark for internal use
        if payload.internalUseOnly:
            try:
                textbox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(3.0), Inches(0.3))
                tf = textbox.text_frame
                tf.text = "Bell Confidential"
                tf.paragraphs[0].font.size = Pt(9)
            except Exception:
                pass  # Skip watermark if it fails

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)

    filename = f"{payload.deck.title or 'Presentation'}.pptx"
    return StreamingResponse(
        output,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={filename}"},
    )

